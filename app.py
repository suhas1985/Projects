import streamlit as st
import os
import PyPDF2
import docx
import pptx
import tempfile
import faiss
import numpy as np
from dotenv import load_dotenv
from langchain.vectorstores import FAISS
from langchain.document_loaders import TextLoader
from langchain.embeddings import OpenAIEmbeddings
from langchain.chains import RetrievalQA
import google.generativeai as palm

# Load environment variables
load_dotenv()

# Function to extract text from PDF using PyPDF2
def extract_text_from_pdf(file):
    text = ""
    with open(file, "rb") as pdf_file:
        pdf_reader = PyPDF2.PdfFileReader(pdf_file)
        for page_num in range(pdf_reader.getNumPages()):
            page = pdf_reader.getPage(page_num)
            text += page.extract_text()
    return text

# Function to extract text from DOCX
def extract_text_from_docx(file):
    doc = docx.Document(file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

# Function to extract text from PPTX
def extract_text_from_pptx(file):
    ppt = pptx.Presentation(file)
    text = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to configure Google Gemini API
def configure_google_gemini(api_key):
    palm.configure(api_key=api_key)

# Function to generate response from Google Gemini API
def generate_response(prompt):
    response = palm.generate_text(
        prompt=prompt,
        max_tokens=150
    )
    return response.result

# Streamlit UI
st.title("RAG Chatbot")
uploaded_file = st.file_uploader("Choose a file", type=["pdf", "docx", "pptx"])

text = ""  # Initialize text variable

if uploaded_file:
    # Save uploaded file to a temporary file
    with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
        tmp_file.write(uploaded_file.read())
        tmp_file_path = tmp_file.name

    # Extract text based on file type
    if uploaded_file.name.endswith(".pdf"):
        text = extract_text_from_pdf(tmp_file_path)
    elif uploaded_file.name.endswith(".docx"):
        text = extract_text_from_docx(tmp_file_path)
    elif uploaded_file.name.endswith(".pptx"):
        text = extract_text_from_pptx(tmp_file_path)
    else:
        st.error("Unsupported file type")
        text = ""

    # Clean up temporary file
    os.remove(tmp_file_path)

    if text:
        # Set up FAISS for retrieval
        loader = TextLoader(text)
        documents = loader.load()
        embeddings = OpenAIEmbeddings()
        vector_store = FAISS.from_documents(documents, embeddings)

        # Set up the retrieval-based QA chain
        qa_chain = RetrievalQA.from_chain_type(
            llm=None,  # We will use Google Gemini API instead of a local LLM
            chain_type="stuff",
            retriever=vector_store.as_retriever()
        )

        # Configure Google Gemini API
        api_key = os.getenv("GOOGLE_API_KEY")
        if not api_key:
            st.error("Google API key not found. Please set it in the .env file.")
        else:
            configure_google_gemini(api_key)
            
            # User prompt for querying the document
            user_prompt = st.text_input("Enter your query about the document:")
            if user_prompt:
                retrieved_docs = qa_chain({"query": user_prompt})["result"]
                response = generate_response(f"{retrieved_docs}\n\nUser: {user_prompt}\nAI:")
                st.write("Response:")
                st.write(response)

# Display the extracted text (for debugging purposes, can be removed)
if text:
    st.subheader("Extracted Text")
    st.write(text)
